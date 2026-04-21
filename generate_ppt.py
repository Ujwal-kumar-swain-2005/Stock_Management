import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(filename="presentation.pptx"):
    prs = Presentation()
    
    # Use the nice abstract background we generated
    bg_image = r"C:\Users\ujwal\.gemini\antigravity\brain\6a3c5fbf-4dfc-41fb-85db-18de0c3609e1\stock_bg_1776773616015.png"

    # We use a blank layout for all slides and manually add text boxes.
    # This ensures that our background image stays correctly at the bottom layer.
    blank_layout = prs.slide_layouts[6]

    def add_slide_with_bg():
        slide = prs.slides.add_slide(blank_layout)
        try:
            slide.shapes.add_picture(bg_image, 0, 0, prs.slide_width, prs.slide_height)
        except Exception as e:
            print("Warning: Could not load background image it will be blank.", e)
        return slide

    def add_title(slide, text, top=Inches(0.5)):
        txBox = slide.shapes.add_textbox(Inches(0.5), top, prs.slide_width - Inches(1), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(255, 255, 255) # White
        p.alignment = PP_ALIGN.LEFT
        return tf

    def add_content(slide, lines, top=Inches(1.8)):
        txBox = slide.shapes.add_textbox(Inches(0.5), top, prs.slide_width - Inches(1), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, (text, level) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.size = Pt(28 - (level * 4))
            p.font.color.rgb = RGBColor(220, 230, 240) # Soft bluish white
        return tf

    # === Slide 1: Title ===
    slide = add_slide_with_bg()
    txBox = slide.shapes.add_textbox(0, Inches(2.5), prs.slide_width, Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Stock Management System"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Project Overview & Architecture\nBy Ujwal Kumar Swain"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 210, 230)
    p2.alignment = PP_ALIGN.CENTER

    # === Slide 2: Introduction ===
    slide = add_slide_with_bg()
    add_title(slide, "What is the Stock Management System?")
    add_content(slide, [
        ("A comprehensive solution to track, organize, and manage business inventory efficiently.", 0),
        ("Key Objectives:", 0),
        ("Automate stock tracking", 1),
        ("Reduce manual errors", 1),
        ("Optimize inventory levels", 1)
    ])

    # === Slide 3: Technologies Used ===
    slide = add_slide_with_bg()
    add_title(slide, "Technologies Built On")
    
    # Left column for frontend
    tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4))
    tf1 = tx1.text_frame
    p = tf1.paragraphs[0]
    p.text = "Frontend"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    for t in ["React (Vite)", "Tailwind CSS", "Axios API client"]:
        p = tf1.add_paragraph()
        p.text = "• " + t
        p.font.size = Pt(24); p.font.color.rgb = RGBColor(220, 230, 240)
        
    # Right column for backend
    tx2 = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4.5), Inches(4))
    tf2 = tx2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Backend & Database"
    p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    for t in ["Spring Boot (Java)", "Spring Security / JWT", "MySQL / PostgreSQL"]:
        p = tf2.add_paragraph()
        p.text = "• " + t
        p.font.size = Pt(24); p.font.color.rgb = RGBColor(220, 230, 240)

    # === Slide 4: Core Features ===
    slide = add_slide_with_bg()
    add_title(slide, "Core Features")
    add_content(slide, [
        ("Dashboard Analytics (Quick insights into stock levels)", 0),
        ("Product Management (Add, edit, delete products)", 0),
        ("Category & Supplier Management", 0),
        ("Order Processing (Tracking inbound/outbound flow)", 0),
        ("Role-Based Access Control (Admin, Manager, Staff)", 0)
    ])

    # === Slide 5: Architecture ===
    slide = add_slide_with_bg()
    add_title(slide, "System Architecture")
    add_content(slide, [
        ("Client (React Frontend)", 0),
        ("RESTful API Communication (JSON)", 1),
        ("Server (Spring Boot Backend)", 0),
        ("Business Logic & Security Validation", 1),
        ("Database (Relational DB)", 0),
        ("Persistent Data Storage", 1)
    ])

    # === Slide 6: Conclusion ===
    slide = add_slide_with_bg()
    add_title(slide, "Conclusion & Future Scope")
    add_content(slide, [
        ("The system provides a solid foundation for robust inventory control.", 0),
        ("Future Enhancements:", 0),
        ("Integration with 3rd-party logistics platforms", 1),
        ("Machine Learning based demand forecasting", 1),
        ("Native mobile application for warehouse staff", 1)
    ])

    prs.save(filename)
    print(f"Presentation saved successfully with custom background as {filename}")

if __name__ == '__main__':
    create_presentation('Stock_Management_System_Presentation_v2.pptx')
